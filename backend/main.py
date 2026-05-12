from dotenv import load_dotenv
load_dotenv()

from fastapi import FastAPI, UploadFile, File, Form, HTTPException
from fastapi.middleware.cors import CORSMiddleware
import google.generativeai as genai
from groq import Groq
import os
from pdf2docx import Converter
import pdfplumber
import pandas as pd
from pptx import Presentation
from pptx.util import Inches
import shutil
import uuid
from utils.rag import RAGManager

rag = RAGManager()

app = FastAPI()

app.add_middleware(
    CORSMiddleware,
    allow_origins=[
        "https://aauca-tools-community.vercel.app",
        "http://localhost:3000",
        "*"
    ],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Logging para depuración
import logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Clientes de IA
GEMINI_KEY = os.getenv("GEMINI_API_KEY")
GROQ_KEY = os.getenv("GROQ_API_KEY")

genai.configure(api_key=GEMINI_KEY)
groq_client = Groq(api_key=GROQ_KEY) if GROQ_KEY else None

# Inicializar modelo Gemini 1.5 Flash para máxima velocidad
gemini_model = genai.GenerativeModel('gemini-1.5-flash')

async def get_ai_response(prompt: str, system_instruction: str = None):
    """Motor de IA híbrido con fallback automático para estabilidad total"""
    # Intentar con Gemini primero
    if GEMINI_KEY:
        try:
            full_prompt = f"{system_instruction}\n\n{prompt}" if system_instruction else prompt
            response = gemini_model.generate_content(full_prompt)
            return response.text
        except Exception as e:
            logger.error(f"Gemini Error, intentando fallback a Groq: {str(e)}")
    
    # Fallback a Groq (Llama 3) si Gemini falla o no hay key
    if GROQ_KEY and groq_client:
        try:
            messages = []
            if system_instruction:
                messages.append({"role": "system", "content": system_instruction})
            messages.append({"role": "user", "content": prompt})
            
            completion = groq_client.chat.completions.create(
                model="llama3-8b-8192", # Modelo ultra rápido
                messages=messages,
                temperature=0.7,
            )
            return completion.choices[0].message.content
        except Exception as e:
            logger.error(f"Groq Fallback Error: {str(e)}")
    
    return "Lo siento, el motor de IA está temporalmente fuera de servicio. Verifica tus API Keys en el servidor."

UPLOAD_DIR = "uploads"
CONVERT_DIR = "converted"
os.makedirs(UPLOAD_DIR, exist_ok=True)
os.makedirs(CONVERT_DIR, exist_ok=True)

@app.post("/convert")
async def convert_file(file: UploadFile = File(...), target_format: str = Form(...)):
    file_id = str(uuid.uuid4())
    input_path = os.path.join(UPLOAD_DIR, f"{file_id}_{file.filename}")
    
    with open(input_path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)
    
    output_filename = f"{file_id}.{target_format}"
    output_path = os.path.join(CONVERT_DIR, output_filename)
    
    try:
        if target_format == "docx":
            cv = Converter(input_path)
            cv.convert(output_path)
            cv.close()
        
        elif target_format == "xlsx":
            with pdfplumber.open(input_path) as pdf:
                all_tables = []
                for page in pdf.pages:
                    table = page.extract_table()
                    if table:
                        all_tables.append(pd.DataFrame(table))
                if all_tables:
                    df = pd.concat(all_tables)
                    df.to_excel(output_path, index=False)
                else:
                    raise HTTPException(status_code=400, detail="No se encontraron tablas en el PDF")
        
        elif target_format == "pptx":
            prs = Presentation()
            with pdfplumber.open(input_path) as pdf:
                for page in pdf.pages:
                    slide = prs.slides.add_slide(prs.slide_layouts[5])
                    text = page.extract_text()
                    if text:
                        textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(5))
                        textbox.text = text
            prs.save(output_path)
            
        return {"download_url": f"/download/{output_filename}", "status": "success"}
    
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

from fastapi.responses import FileResponse

@app.get("/download/{filename}")
async def download_file(filename: str):
    file_path = os.path.join(CONVERT_DIR, filename)
    if os.path.exists(file_path):
        return FileResponse(file_path)
    raise HTTPException(status_code=404, detail="Archivo no encontrado")

@app.post("/transcribe")
async def transcribe(file: UploadFile = File(...)):
    if not GROQ_KEY or not groq_client:
         raise HTTPException(status_code=500, detail="Groq API Key no configurada para transcripción.")

    temp_path = f"temp_{file.filename}"
    with open(temp_path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)
    
    try:
        with open(temp_path, "rb") as audio_file:
            transcription = groq_client.audio.transcriptions.create(
                file=(temp_path, audio_file.read()),
                model="whisper-large-v3",
                response_format="text"
            )
        return {"text": transcription}
    finally:
        os.remove(temp_path)

@app.post("/summarize")
async def summarize(text: str = Form(...)):
    instruction = "Eres un experto académico de la AAUCA. Resume el contenido de forma estructurada, usando negritas y puntos clave."
    response = await get_ai_response(f"Resume esto:\n\n{text}", instruction)
    return {"summary": response}

@app.post("/upload-pdf")
async def upload_pdf(file: UploadFile = File(...)):
    file_id = str(uuid.uuid4())
    input_path = os.path.join(UPLOAD_DIR, f"{file_id}_{file.filename}")
    
    with open(input_path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)
    
    try:
        with pdfplumber.open(input_path) as pdf:
            for i, page in enumerate(pdf.pages):
                text = page.extract_text()
                if text:
                    await rag.store_chunk(file_id, text, i + 1)
        
        return {"status": "success", "file_id": file_id}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/chat")
async def chat_with_docs(question: str = Form(...), file_id: str = Form(None)):
    try:
        context_chunks = await rag.query(question, file_id)
        context_text = "\n".join([c['content'] for c in context_chunks]) if context_chunks else "Sin contexto adicional del documento."
        
        instruction = f"Eres el Tutor de IA de Élite de la AAUCA, desarrollado por Tranquilino Mba Ncogo. Ayuda al estudiante usando este contexto:\n\n{context_text}"
        response = await get_ai_response(question, instruction)
        
        return {"response": response}
    except Exception as e:
        logger.error(f"Chat error: {str(e)}")
        raise HTTPException(status_code=500, detail="Error crítico en el motor de chat.")

@app.post("/clear-cache")
async def clear_cache(file_id: str = Form(None)):
    """Limpieza de base de datos y archivos locales para optimizar espacio"""
    try:
        if file_id:
            await rag.delete_document(file_id)
        
        # Limpiar archivos temporales de más de 1 hora
        import time
        now = time.time()
        for folder in [UPLOAD_DIR, CONVERT_DIR]:
            for f in os.listdir(folder):
                f_path = os.path.join(folder, f)
                if os.stat(f_path).st_mtime < now - 3600:
                    if os.path.isfile(f_path):
                        os.remove(f_path)
        
        return {"status": "success", "message": "Cache y base de datos optimizados"}
    except Exception as e:
        logger.error(f"Clear cache error: {str(e)}")
        return {"status": "error", "message": str(e)}

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000)
